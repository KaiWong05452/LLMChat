# To extract the texts from common documents including ppt, pdf, doc.
import os
import pypdf
import pptx
import docx


# Function to extract text from a file based on its extension
def extract_text_from_file(file_path):
    # Get the file extension
    file_extension = os.path.splitext(file_path)[1]
    # Check the file extension and call the appropriate function
    if file_extension == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension == '.pptx':
        return extract_text_from_pptx(file_path)
    elif file_extension == '.docx':
        return extract_text_from_docx(file_path)
    elif file_extension == '.txt':
        return extract_text_from_txt(file_path)
    else:
        # If the file extension is not supported, return None
        return None


# Function to extract text from a txt file
def extract_text_from_txt(file_path):
    # Open the txt file
    with open(file_path, 'r', encoding='utf-8') as file:
        # Read the contents of the file
        text = file.read()
    # Return the extracted text
    return text


# Function to extract text from a PDF file
def extract_text_from_pdf(file_path):
    # Open the PDF file
    pdf = pypdf.PdfReader(file_path)
    text = ''
    # Loop through each page in the PDF
    for page in pdf.pages:
        # Extract the text from the page and append it to the text variable
        text += page.extract_text()
    # Return the extracted text
    return text


# Function to extract text from a DOCX file
def extract_text_from_docx(file_path):
    # Open the DOCX file
    doc = docx.Document(file_path)
    text = ''
    # Loop through each paragraph in the DOCX file
    for paragraph in doc.paragraphs:
        # Append the text from the paragraph to the text variable
        text += paragraph.text
    # Return the extracted text
    return text


# Function to extract text from a PPTX file
def extract_text_from_pptx(file_path):
    # Open the PPTX file
    pres = pptx.Presentation(file_path)
    text = ''
    # Loop through each slide in the PPTX file
    for slide in pres.slides:
        # Loop through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has a text frame
            if shape.has_text_frame:
                # Loop through each paragraph in the text frame
                for paragraph in shape.text_frame.paragraphs:
                    # Loop through each run in the paragraph
                    for run in paragraph.runs:
                        # Append the text from the run to the text variable
                        text += run.text
    # Return the extracted text
    return text


def save_file_and_extract(file, output_dir):
    if file is not None:
        file_path = os.path.join(output_dir, file.filename)
        file.save(file_path)

        document = extract_text_from_file(file_path)
    else:
        document = None

    return document
